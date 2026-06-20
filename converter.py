import os
import io
import time
import zipfile
import pandas as pd
import PyPDF2
import docx
from pptx import Presentation

class DocumentConverter:
    """Class to parse different document types to Markdown format."""
    
    
    @staticmethod
    def convert_pdf(file_bytes, progress_cb):
        """Extract text page by page from PDF files."""
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(file_bytes))
        total_pages = len(pdf_reader.pages)
        markdown_content = ""
        
        for idx in range(total_pages):
            progress_cb(f"Extracting text from page {idx + 1} of {total_pages}...")
            page = pdf_reader.pages[idx]
            text = page.extract_text()
            if text:
                markdown_content += f"## Page {idx + 1}\n\n{text}\n\n"
        return markdown_content

    @staticmethod
    def convert_docx(file_bytes, progress_cb):
        """Parse Word DOCX document hierarchy including paragraph styles and tables."""
        from docx import Document
        from docx.document import Document as _Document
        from docx.table import Table, _Cell
        from docx.text.paragraph import Paragraph
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P

        def iter_block_items(parent):
            """Yield paragraph and table elements in visual document order."""
            if isinstance(parent, _Document):
                parent_elm = parent.element.body
            elif isinstance(parent, _Cell):
                parent_elm = parent._tc
            else:
                raise ValueError("Unsupported docx parent type")

            for child in parent_elm.iterchildren():
                if isinstance(child, CT_P):
                    yield Paragraph(child, parent)
                elif isinstance(child, CT_Tbl):
                    yield Table(child, parent)

        doc = Document(io.BytesIO(file_bytes))
        markdown_content = ""
        
        blocks = list(iter_block_items(doc))
        total_blocks = len(blocks)
        
        for idx, block in enumerate(blocks):
            if idx % 10 == 0:
                progress_cb(f"Processing content structure... element {idx + 1} of {total_blocks}")
                
            if isinstance(block, Paragraph):
                text = block.text.strip()
                if not text:
                    continue
                
                # Map style names to markdown structures
                style_name = block.style.name.lower()
                if "heading 1" in style_name:
                    markdown_content += f"# {text}\n\n"
                elif "heading 2" in style_name:
                    markdown_content += f"## {text}\n\n"
                elif "heading 3" in style_name:
                    markdown_content += f"### {text}\n\n"
                elif "heading 4" in style_name:
                    markdown_content += f"#### {text}\n\n"
                elif "list bullet" in style_name or "bullet" in style_name:
                    markdown_content += f"- {text}\n"
                elif "list number" in style_name or "number" in style_name:
                    markdown_content += f"1. {text}\n"
                else:
                    markdown_content += f"{text}\n\n"
            elif isinstance(block, Table):
                # Standard markdown table converter
                md_rows = []
                for row_idx, row in enumerate(block.rows):
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    md_rows.append("| " + " | ".join(cells) + " |")
                    if row_idx == 0:
                        md_rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                markdown_content += "\n".join(md_rows) + "\n\n"
                
        return markdown_content

    @staticmethod
    def convert_pptx(file_bytes, progress_cb):
        """Parse PowerPoint PPTX slides and create clean outlining structures."""
        prs = Presentation(io.BytesIO(file_bytes))
        markdown_content = ""
        total_slides = len(prs.slides)
        
        for idx, slide in enumerate(prs.slides, 1):
            progress_cb(f"Parsing slides... slide {idx} of {total_slides}")
            markdown_content += f"## Slide {idx}\n\n"
            
            # Find and highlight slide title
            title_text = ""
            if slide.shapes.title:
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    markdown_content += f"### {title_text}\n\n"
                    
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            # Apply indentation level mapping
                            level = paragraph.level
                            if level > 0:
                                indent = "  " * level
                                markdown_content += f"{indent}- {text}\n"
                            else:
                                markdown_content += f"{text}\n\n"
            markdown_content += "\n---\n\n"
        return markdown_content

    @staticmethod
    def convert_csv(file_bytes, progress_cb):
        """Convert CSV dataset to a neat markdown table."""
        progress_cb("Importing dataset elements...")
        df = pd.read_csv(io.BytesIO(file_bytes))
        return df.to_markdown(index=False) + "\n\n"

    @staticmethod
    def convert_excel(file_bytes, progress_cb):
        """Convert Excel spreadsheets (all sheets) to markdown tables."""
        progress_cb("Scanning spreadsheets sheets...")
        excel_data = pd.read_excel(io.BytesIO(file_bytes), sheet_name=None)
        markdown_content = ""
        for sheet_name, df in excel_data.items():
            markdown_content += f"## Sheet: {sheet_name}\n\n"
            markdown_content += df.to_markdown(index=False) + "\n\n"
        return markdown_content

    @staticmethod
    def convert_txt(file_bytes, progress_cb):
        """Safely read and output raw text files."""
        progress_cb("Reading text contents...")
        return file_bytes.decode("utf-8", errors="replace")


def convert_single_file(ext, file_bytes, filename, progress_cb):
    """Router to send file bytes to the appropriate format parser."""
    if ext == ".pdf":
        return DocumentConverter.convert_pdf(file_bytes, progress_cb)
    elif ext == ".docx":
        return DocumentConverter.convert_docx(file_bytes, progress_cb)
    elif ext == ".pptx":
        return DocumentConverter.convert_pptx(file_bytes, progress_cb)
    elif ext == ".csv":
        return DocumentConverter.convert_csv(file_bytes, progress_cb)
    elif ext in (".xlsx", ".xls"):
        return DocumentConverter.convert_excel(file_bytes, progress_cb)
    elif ext == ".txt":
        return DocumentConverter.convert_txt(file_bytes, progress_cb)
    else:
        raise ValueError(f"Extension not supported: {ext}")


def process_zip_item(zip_bytes, results_dict, zip_out, update_filename_cb, update_detail_cb):
    """Extract and recursively convert all supported document formats in a ZIP archive."""
    supported = (".pdf", ".docx", ".pptx", ".csv", ".xlsx", ".xls", ".txt")
    with zipfile.ZipFile(io.BytesIO(zip_bytes), "r") as zip_in:
        namelist = zip_in.namelist()
        valid_items = [n for n in namelist if not n.endswith("/") and os.path.splitext(n)[1].lower() in supported]
        total_items = len(valid_items)
        
        for idx, name in enumerate(valid_items, 1):
            ext = os.path.splitext(name)[1].lower()
            update_filename_cb(f"{name} (Archive File {idx}/{total_items})")
            try:
                file_bytes = zip_in.read(name)
                # Parse using the router
                md_content = convert_single_file(ext, file_bytes, name, update_detail_cb)
                
                # Write to out zip
                out_name = os.path.splitext(name)[0] + ".md"
                zip_out.writestr(out_name, md_content)
                
                results_dict[out_name] = {
                    "name": name,
                    "content": md_content,
                    "status": "Success",
                    "size_kb": len(md_content) / 1024
                }
            except Exception as e:
                results_dict[name] = {
                    "name": name,
                    "content": f"# Error converting {name}\n\n```{str(e)}```",
                    "status": f"Failed: {str(e)}",
                    "size_kb": 0
                }


def run_conversion_pipeline(files_to_process, source_mode="browser", update_ui_cb=None):
    """Unified file conversion pipeline executing independently from Streamlit context."""
    zip_buffer = io.BytesIO()
    results = {}
    
    start_time = time.time()
    total_files = len(files_to_process)
    
    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_out:
        for idx, item in enumerate(files_to_process, 1):
            filename = item["name"]
            is_zip = item["is_zip"]
            
            # Interactive callback for UI
            def progress_cb(detail):
                if update_ui_cb:
                    update_ui_cb(filename, idx, total_files, detail)
            
            if is_zip:
                if source_mode == "browser":
                    file_bytes = item["bytes"]
                else:
                    with open(item["full_path"], "rb") as f:
                        file_bytes = f.read()
                
                def update_zip_fname(sub_name):
                    if update_ui_cb:
                        update_ui_cb(sub_name, idx, total_files, f"Extracting ZIP Archive: {filename}")
                
                process_zip_item(
                    file_bytes, 
                    results, 
                    zip_out, 
                    update_zip_fname, 
                    progress_cb
                )
            else:
                ext = os.path.splitext(filename)[1].lower()
                if update_ui_cb:
                    update_ui_cb(filename, idx, total_files, "Reading file data...")
                
                try:
                    if source_mode == "browser":
                        file_bytes = item["bytes"]
                    else:
                        with open(item["full_path"], "rb") as f:
                            file_bytes = f.read()
                    
                    md_content = convert_single_file(ext, file_bytes, filename, progress_cb)
                    
                    out_name = os.path.splitext(filename)[0] + ".md"
                    zip_out.writestr(out_name, md_content)
                    
                    results[out_name] = {
                        "name": filename,
                        "content": md_content,
                        "status": "Success",
                        "size_kb": len(md_content) / 1024
                    }
                except Exception as e:
                    results[filename] = {
                        "name": filename,
                        "content": f"# Error converting {filename}\n\n```{str(e)}```",
                        "status": f"Failed: {str(e)}",
                        "size_kb": 0
                    }
                    
    duration = time.time() - start_time
    return zip_buffer.getvalue(), results, duration
